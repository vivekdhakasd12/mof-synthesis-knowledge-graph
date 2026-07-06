"""Build the Case Study 2 exposé presentation (.pptx).

Run with uv so python-pptx does not need a global install:
    uv run --with python-pptx python presentation/build_slides.py

Output: presentation/Case_Study_2_Expose_Presentation.pptx

Content comes only from the exposé report (docs/expose.pdf). The visual style
matches docs/expose.html: SRH orange #E64415 accent, Lato font, srh_logo.jpg,
clean white, charcoal text, thin gray rules, minimalist.

Writing rules: no em dashes. En dash only in number/date ranges and in the term
"Metal-Organic Framework" (rendered with an en dash to match the exposé).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---- SRH exposé palette (from docs/expose.html) ----------------------------
ORANGE = RGBColor(0xE6, 0x44, 0x15)   # SRH accent
INK = RGBColor(0x1A, 0x1A, 0x1A)      # headings / body
GRAY = RGBColor(0x55, 0x55, 0x55)     # secondary / captions
RULE = RGBColor(0xE3, 0xE3, 0xE3)     # thin section rules
RULE2 = RGBColor(0xBB, 0xBB, 0xBB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0xFA, 0xF6, 0xF4)    # faint warm tint for callout box

FONT = "Lato"

DOCS = Path("/Users/dev/Agentic Workflows /case-study-2/docs")
LOGO = str(DOCS / "assets" / "srh_logo.jpg")
LOGO_RATIO = 457 / 591  # height / width

# en dash and >= as used in the exposé
ND = "–"
GE = "≥"

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---- helpers ---------------------------------------------------------------
def _no_shadow(shape):
    # Flat look: the default shadow comes from the shape's <p:style> effectRef.
    # Remove that style and add an explicit empty effect list.
    el = shape._element
    style = el.find(qn("p:style"))
    if style is not None:
        el.remove(style)
    spPr = el.spPr
    if spPr.find(qn("a:effectLst")) is None:
        spPr.append(
            parse_xml('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        )


def slide():
    return prs.slides.add_slide(BLANK)


def R(txt, size=14, color=INK, bold=False, italic=False):
    return (txt, size, color, bold, italic)


def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.12):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic) in para:
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = FONT
    return tb


def rule(s, x, y, w, color=RULE, pt=1.0):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(pt))
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    _no_shadow(r)
    return r


def box(s, x, y, w, h, fill=None, line=None, pt=1.0):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(pt)
    b.shadow.inherit = False
    _no_shadow(b)
    return b


def logo_small(s):
    w = Inches(0.95)
    s.shapes.add_picture(LOGO, Inches(0.6), Inches(0.46), width=w, height=Emu(int(w * LOGO_RATIO)))


def pagenum(s, n):
    text(s, Inches(12.5), Inches(7.04), Inches(0.7), Inches(0.3),
         [[R(str(n), 10, RULE2)]], align=PP_ALIGN.RIGHT)


def header(s, num, title, n):
    """Logo, orange number + charcoal title, thin rule. Returns body top (Inches)."""
    logo_small(s)
    runs = []
    if num:
        runs.append(R(num + "   ", 27, ORANGE, bold=True))
    runs.append(R(title, 27, INK, bold=True))
    text(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.7), [runs])
    rule(s, Inches(0.64), Inches(2.18), Inches(12.05), color=RULE, pt=1.2)
    pagenum(s, n)
    return Inches(2.5)


def bullets(s, x, y, w, items, size=15, gap=9, lead=1.16):
    """items: list of (text, level). level 0 = orange bullet, 1 = indented dash."""
    paras = []
    for txt, lvl in items:
        if lvl == 0:
            paras.append([R("•  ", size, ORANGE, bold=True), R(txt, size, INK)])
        else:
            paras.append([R("      –  ", size - 1, ORANGE), R(txt, size - 1, GRAY)])
    return text(s, x, y, w, Inches(4.4), paras, space_after=gap, line_spacing=lead)


# ===========================================================================
# Slide 1: Title (mirrors the exposé title page)
# ===========================================================================
s = slide()
lw = Inches(1.7)
s.shapes.add_picture(LOGO, Emu(int((SW - lw) / 2)), Inches(1.15), width=lw, height=Emu(int(lw * LOGO_RATIO)))
text(s, Inches(1.0), Inches(2.85), Inches(11.33), Inches(0.4),
     [[R("C A S E   S T U D Y   2      ·      E X P O S É", 14, GRAY, bold=True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(0.9), Inches(3.45), Inches(11.53), Inches(1.5),
     [[R("Building and Validating a Metal" + ND + "Organic Framework", 26, INK, bold=True)],
      [R("Synthesis Knowledge Graph with Large Language Models", 26, INK, bold=True)]],
     align=PP_ALIGN.CENTER, line_spacing=1.08)
rule(s, Emu(int((SW - Inches(1.4)) / 2)), Inches(4.95), Inches(1.4), color=ORANGE, pt=3.0)
text(s, Inches(1.0), Inches(5.45), Inches(11.33), Inches(1.5),
     [[R("Devendra Singh Dhakad", 16, INK, bold=True),
       R("      ·      Matriculation No. 100004684", 16, GRAY)],
      [R("M.Sc. Data Science and Artificial Intelligence", 14, GRAY)],
      [R("Supervisor: Prof. Dr. Mehrdad Jalali", 14, GRAY)],
      [R("SRH University of Applied Sciences Heidelberg", 14, GRAY)]],
     align=PP_ALIGN.CENTER, space_after=7)
text(s, Inches(1.0), Inches(6.9), Inches(11.33), Inches(0.35),
     [[R("11 June 2026", 12, RULE2)]], align=PP_ALIGN.CENTER)


# ===========================================================================
# Slide 2: Introduction and relevance
# ===========================================================================
s = slide()
y = header(s, "2", "Introduction and relevance", 2)
bullets(s, Inches(0.7), y, Inches(11.9), [
    ("Metal" + ND + "Organic Frameworks (MOFs) are among the most intensively studied "
     "porous crystalline materials: over 100,000 synthesised structures are recorded "
     "in the Cambridge Structural Database.", 0),
    ("Their applications span CO2 capture, gas storage and separation, catalysis, "
     "sensing, and drug delivery.", 0),
    ("Yet how a MOF is made (precursor, linker, solvent, method, temperature, time) "
     "stays locked in the prose of tens of thousands of publications.", 0),
    ("So synthesis planning still relies heavily on expert intuition and trial and error.", 0),
], size=16, gap=14)


# ===========================================================================
# Slide 3: The opportunity and the gap
# ===========================================================================
s = slide()
y = header(s, "", "The opportunity, and the gap", 3)
bullets(s, Inches(0.7), y, Inches(11.9), [
    ("Large Language Models and Knowledge Graphs can turn this text into a structured, "
     "queryable resource (Pan et al., 2024).", 0),
    ("Bai et al. (2025) built a 2.5-million-node knowledge graph for framework materials "
     "from over 100,000 articles. Scale is already proven.", 0),
    ("What stays under-quantified is reliability:", 0),
    ("How accurate are LLM-extracted synthesis records, field by field?", 1),
    ("How do they compare with established text-mined databases and domain baselines?", 1),
    ("What do open-weight models deliver versus commercial APIs, and at what cost?", 1),
], size=16, gap=11)


# ===========================================================================
# Slide 4: This project
# ===========================================================================
s = slide()
y = header(s, "", "This project: build and validate", 4)
bullets(s, Inches(0.7), y, Inches(11.9), [
    ("An end-to-end pipeline extracts complete synthesis records (precursors, linkers, "
     "solvents, methods, conditions, properties, applications) from open-access MOF "
     "literature into a Neo4j graph with full provenance.", 0),
    ("Extraction quality is validated three ways:", 0),
    ("against a hand-annotated gold standard,", 1),
    ("cross-checked against the DigiMOF and SynMOF databases,", 1),
    ("benchmarked against ChemDataExtractor and MatSciBERT.", 1),
    ("The outcome is a reproducible validation framework for LLM-based scientific "
     "extraction, plus a queryable, provenance-aware MOF synthesis knowledge graph.", 0),
], size=16, gap=11)


# ===========================================================================
# Slide 5: Research question
# ===========================================================================
s = slide()
y = header(s, "3", "Research question", 5)
box(s, Inches(0.7), y, Inches(11.9), Inches(1.25), fill=FAINT)
box(s, Inches(0.7), y, Inches(0.1), Inches(1.25), fill=ORANGE)
text(s, Inches(1.0), y + Inches(0.16), Inches(11.4), Inches(1.0),
     [[R("Main question:  ", 14.5, ORANGE, bold=True),
       R("How accurately and reliably can LLMs extract complete MOF synthesis records "
         "(precursor, linker, solvent, method, conditions) from the literature, measured "
         "against expert annotation, established text-mined databases (DigiMOF, SynMOF), "
         "and domain-specific baselines?", 14.5, INK)]], line_spacing=1.12)
bullets(s, Inches(0.7), y + Inches(1.5), Inches(11.9), [
    ("Which prompting strategy (zero-shot, few-shot, schema-guided, chain-of-thought) "
     "is most reliable for each field?", 0),
    ("Where and why do LLM extractions disagree with DigiMOF and SynMOF, and which "
     "source is correct?", 0),
    ("How do open-weight models (Llama-3) compare to commercial APIs (GPT-4o, Claude) "
     "on per-field accuracy, cost, and latency?", 0),
    ("Can the KG answer cross-paper aggregation queries, and, as a stretch goal, support "
     "KG-grounded question answering?", 0),
], size=14, gap=7)


# ===========================================================================
# Slide 6: Objectives
# ===========================================================================
s = slide()
y = header(s, "", "Objectives", 6)
bullets(s, Inches(0.7), y, Inches(11.9), [
    ("Finalise the MOF synthesis ontology (v0.2: eight entity types plus mandatory "
     "provenance) with supervisor input.", 0),
    ("Assemble 300" + ND + "500 open-access MOF synthesis papers (CSD MOF subset, the "
     "DigiMOF article index, ChemRxiv, and PubMed Central OA).", 0),
    ("Implement a unified extraction interface spanning LLM strategies and baselines "
     "(ChemDataExtractor 2.0, MatSciBERT).", 0),
    ("Hand-annotate 150" + ND + "200 synthesis paragraphs in Label Studio and define the "
     "DigiMOF/SynMOF field mapping.", 0),
    ("Construct a Neo4j KG with chemical-name entity resolution and a Streamlit dashboard "
     "with natural-language-to-Cypher querying.", 0),
], size=15.5, gap=12)


# ===========================================================================
# Slide 7: State of the art
# ===========================================================================
s = slide()
y = header(s, "4", "State of the art", 7)
themes = [
    ("Rule-based text mining: ", "ChemDataExtractor (Swain & Cole, 2016) powers DigiMOF "
     "(Glasby, 2023) and SynMOF (Luo, 2022). Precise but rigid on free-text conditions."),
    ("Domain language models: ", "MatSciBERT (Gupta, 2022) improves materials NER but "
     "needs labelled data and stops at entity level, not complete records."),
    ("LLM extraction in chemistry: ", "Zheng (2023), Dagdelen (2024), Polak & Morgan "
     "(2024), Shi (2024), Lin (2025). Mostly commercial GPT, validated on small sets."),
    ("LLM-constructed knowledge graphs: ", "Pan (2024), Bai (2025), GraphRAG (Edge, 2024). "
     "Scale is shown; per-field extraction validity is not the focus."),
    ("Graph ML on MOF data: ", "MOFGalaxyNet (Jalali, 2023) predicts properties from a "
     "MOF network, a downstream use that depends on reliable structured data."),
]
paras = []
for head, body in themes:
    paras.append([R("•  ", 13.5, ORANGE, bold=True), R(head, 13.5, INK, bold=True),
                  R(body, 13.5, GRAY)])
text(s, Inches(0.7), y, Inches(11.9), Inches(3.2), paras, space_after=9, line_spacing=1.12)
box(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.85), fill=FAINT)
box(s, Inches(0.7), Inches(6.05), Inches(0.1), Inches(0.85), fill=ORANGE)
text(s, Inches(1.0), Inches(6.16), Inches(11.4), Inches(0.7),
     [[R("The gap:  ", 13.5, ORANGE, bold=True),
       R("no study validates LLM-extracted MOF synthesis records field by field against "
         "both expert annotation and the text-mined databases. This project closes it.",
         13.5, INK)]], line_spacing=1.1)


# ===========================================================================
# Slide 8: Methodology (5 phases)
# ===========================================================================
s = slide()
y = header(s, "5", "Methodology", 8)
phases = [
    ("1", "Corpus collection", "300" + ND + "500 open-access papers (CSD MOF subset, "
     "DigiMOF index, ChemRxiv, PubMed Central OA), parsed to clean text with GROBID / PyMuPDF."),
    ("2", "Schema and annotation", "MOF ontology v0.2; 150" + ND + "200 paragraphs "
     "hand-annotated in Label Studio as the gold standard; field mapping to DigiMOF/SynMOF."),
    ("3", "Extraction pipelines", "LLMs (GPT-4o, Claude, Llama-3) under zero-shot, few-shot, "
     "schema-guided and chain-of-thought prompting; baselines ChemDataExtractor 2.0 and "
     "MatSciBERT; one shared output schema."),
    ("4", "KG construction", "Neo4j with SBERT chemical-name entity resolution and "
     "provenance edges; target 5,000+ entities and 25,000+ relations."),
    ("5", "Validation and dashboard", "Per-field precision / recall / F1 versus the gold "
     "standard; agreement versus DigiMOF/SynMOF; cost and latency profiling; Streamlit "
     "natural-language-to-Cypher dashboard."),
]
ry = y
for num, name, body in phases:
    text(s, Inches(0.7), ry, Inches(0.55), Inches(0.5), [[R(num, 18, ORANGE, bold=True)]])
    text(s, Inches(1.3), ry, Inches(3.0), Inches(0.9), [[R(name, 14, INK, bold=True)]],
         anchor=MSO_ANCHOR.TOP, line_spacing=1.05)
    text(s, Inches(4.4), ry, Inches(8.2), Inches(0.9), [[R(body, 13, GRAY)]], line_spacing=1.1)
    ry = ry + Inches(0.86)
    if num != "5":
        rule(s, Inches(0.7), ry - Inches(0.12), Inches(11.9), color=RULE, pt=0.75)


# ===========================================================================
# Slide 9: Work plan and timeline
# ===========================================================================
s = slide()
y = header(s, "6", "Work plan and timeline", 9)
text(s, Inches(0.7), y - Inches(0.05), Inches(11.9), Inches(0.4),
     [[R("12 weeks  (22 June " + ND + " 15 September 2026)", 13, GRAY, italic=True)]])
weeks = [
    ("1" + ND + "2", "Literature deep dive; finalise MOF ontology with supervisor; assemble corpus DOI lists"),
    ("3", "PDF parsing pipeline; cleaned corpus; synthesis-paragraph coverage analysis"),
    ("4" + ND + "5", "Implement LLM extraction pipelines (zero-shot, few-shot, schema-guided, CoT) and baselines"),
    ("6", "Hand-annotate the gold standard; freeze the evaluation protocol and DigiMOF/SynMOF field mapping"),
    ("7", "Set up Neo4j; KG ingestion; entity resolution and chemical-name normalisation"),
    ("8", "Run all extractors at scale; populate the full KG; compute per-field metrics"),
    ("9", "Agreement analysis versus DigiMOF/SynMOF; error taxonomy; cost and latency comparison; ablations"),
    ("10", "Streamlit dashboard with natural-language-to-Cypher and provenance views; stretch: KG-RAG QA"),
    ("11", "End-to-end reproducibility (fresh Docker build); supervisor feedback round; polish"),
    ("12", "Finalise the report, code repository, and presentation slides"),
]
ty = y + Inches(0.5)
for wk, task in weeks:
    text(s, Inches(0.7), ty, Inches(1.05), Inches(0.4), [[R(wk, 12.5, ORANGE, bold=True)]])
    text(s, Inches(1.75), ty, Inches(10.85), Inches(0.45), [[R(task, 12.5, INK)]], line_spacing=1.05)
    ty = ty + Inches(0.42)


# ===========================================================================
# Slide 10: Expected results
# ===========================================================================
s = slide()
y = header(s, "7", "Expected results", 10)
box(s, Inches(0.7), y, Inches(5.8), Inches(3.5), fill=FAINT)
box(s, Inches(0.7), y, Inches(0.1), Inches(3.5), fill=ORANGE)
text(s, Inches(1.0), y + Inches(0.25), Inches(5.3), Inches(3.1),
     [[R("Quantitative", 16, ORANGE, bold=True)],
      [R("Per-field precision, recall and F1 for every extractor and prompting "
         "combination (target F1 " + GE + " 0.80 on well-defined fields such as precursor, "
         "linker and solvent; honest reporting on harder fields such as process conditions).",
         13.5, INK)],
      [R("Agreement rates with DigiMOF and SynMOF on overlapping papers.", 13.5, INK)],
      [R("A cost and latency table for commercial versus open-weight models.", 13.5, INK)]],
     space_after=8, line_spacing=1.12)
box(s, Inches(6.85), y, Inches(5.75), Inches(3.5), fill=FAINT)
box(s, Inches(6.85), y, Inches(0.1), Inches(3.5), fill=ORANGE)
text(s, Inches(7.15), y + Inches(0.25), Inches(5.25), Inches(3.1),
     [[R("Qualitative", 16, ORANGE, bold=True)],
      [R("A failure-mode taxonomy of where LLM extraction breaks:", 13.5, INK)],
      [R("hallucinated values, unit errors, missed coreferences, and schema violations.",
         13.5, GRAY)],
      [R("Per-field guidance on which prompting strategy to trust for each kind of "
         "synthesis field.", 13.5, INK)]],
     space_after=8, line_spacing=1.12)


# ===========================================================================
# Slide 11: Artefact and deliverables
# ===========================================================================
s = slide()
y = header(s, "", "Artefact and deliverables", 11)
text(s, Inches(0.7), y, Inches(11.9), Inches(0.7),
     [[R("Artefact:  ", 15, ORANGE, bold=True),
       R("a provenance-aware Neo4j MOF synthesis knowledge graph (target 5,000+ entities, "
         "25,000+ relations), queryable via a Streamlit dashboard that translates "
         "natural-language questions to Cypher.", 15, INK)]], line_spacing=1.15)
bullets(s, Inches(0.7), y + Inches(1.15), Inches(11.9), [
    ("Reproducible Docker-packaged repository.", 0),
    ("KG dump, ontology, and gold standard.", 0),
    ("Evaluation report with full metrics and error analysis.", 0),
    ("Final written report and defence slides.", 0),
], size=14.5, gap=7)
box(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.15), fill=FAINT)
box(s, Inches(0.7), Inches(5.7), Inches(0.1), Inches(1.15), fill=ORANGE)
text(s, Inches(1.0), Inches(5.84), Inches(11.4), Inches(0.95),
     [[R("Ethics and reproducibility:  ", 13.5, ORANGE, bold=True),
       R("only open-access publications under permissive licences are processed. The "
         "pipeline is containerised for one-command reproducibility, and all LLM calls "
         "are cached and logged so experiments replay without re-incurring inference cost.",
         13.5, INK)]], line_spacing=1.12)


# ===========================================================================
# Slide 12: Outlook
# ===========================================================================
s = slide()
y = header(s, "", "Outlook", 12)
bullets(s, Inches(0.7), y, Inches(11.9), [
    ("A validated extraction pipeline turns literature into data: the KG can feed "
     "synthesis-prediction models (Luo et al., 2022) and graph-based property models "
     "such as MOFGalaxyNet (Jalali et al., 2023).", 0),
    ("The framework extends naturally to covalent organic frameworks, a direct path to "
     "follow-on thesis work.", 0),
], size=16, gap=14)
rule(s, Inches(0.7), Inches(5.5), Inches(11.9), color=RULE, pt=1.2)
text(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.9),
     [[R("Devendra Singh Dhakad", 15, INK, bold=True),
       R("   ·   Matriculation No. 100004684   ·   ", 13, GRAY),
       R("SRH University of Applied Sciences Heidelberg", 13, GRAY)],
      [R("Supervisor: Prof. Dr. Mehrdad Jalali", 12.5, GRAY)]], space_after=6)


# ===========================================================================
# Slide 13: References (appendix)
# ===========================================================================
s = slide()
y = header(s, "", "References", 13)
refs = [
    "Bai, X., et al. (2025). Construction of a Knowledge Graph for Framework Material Enabled by LLMs. npj Comput. Mater., 11, 51.",
    "Dagdelen, J., et al. (2024). Structured Information Extraction from Scientific Text with LLMs. Nat. Commun., 15, 1418.",
    "Edge, D., et al. (2024). From Local to Global: A GraphRAG Approach to Query-Focused Summarization. arXiv:2404.16130.",
    "Glasby, L. T., et al. (2023). DigiMOF: A Database of MOF Synthesis Information. Chem. Mater., 35(11), 4510" + ND + "4524.",
    "Gupta, T., et al. (2022). MatSciBERT: A Materials Domain Language Model. npj Comput. Mater., 8, 102.",
    "Jalali, M., et al. (2023). MOFGalaxyNet: Social Network Analysis for Predicting Guest Accessibility in MOFs. J. Cheminform., 15, 94.",
    "Lin, Z., et al. (2025). Reshaping MOFs Text Mining with a Dynamic Multi-Agent Framework. arXiv:2504.18880.",
    "Luo, Y., et al. (2022). MOF Synthesis Prediction Enabled by Automatic Data Mining and ML. Angew. Chem. Int. Ed., 61, e202200242.",
    "Pan, S., et al. (2024). Unifying Large Language Models and Knowledge Graphs: A Roadmap. IEEE TKDE, 36(7), 3580" + ND + "3599.",
    "Polak, M. P., & Morgan, D. (2024). Extracting Accurate Materials Data with Conversational LLMs. Nat. Commun., 15, 1569.",
    "Shi, L., et al. (2024). LLM-Based MOFs Synthesis Condition Extraction Using Few-Shot Demonstrations. arXiv:2408.04665.",
    "Swain, M. C., & Cole, J. M. (2016). ChemDataExtractor. J. Chem. Inf. Model., 56(10), 1894" + ND + "1904.",
    "Zheng, Z., et al. (2023). ChatGPT Chemistry Assistant for Text Mining and Prediction of MOF Synthesis. JACS, 145(32), 18048" + ND + "18062.",
]
paras = [[R("•  ", 10.5, ORANGE), R(r, 10.5, GRAY)] for r in refs]
text(s, Inches(0.7), y - Inches(0.05), Inches(11.9), Inches(4.6), paras,
     space_after=4, line_spacing=1.06)


out = Path(__file__).resolve().parent / "Case_Study_2_Expose_Presentation.pptx"
prs.save(str(out))
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
